from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from core.registry import register
from core.schemas import ParsedDoc, TextBlock, ImageAsset
from core.parsing.base import BaseParser


@register("parsing", "pptx")
class PptxParser(BaseParser):
    def parse(self, path: Path, doc_id: str) -> ParsedDoc:
        prs = Presentation(path)
        text_blocks: list[TextBlock] = []
        images: list[ImageAsset] = []
        img_n = 0

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    slide_text.append(shape.text_frame.text.strip())
                if self.extract_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_n += 1
                    image = shape.image
                    img_id = f"{doc_id}_img{img_n}"
                    out_path = self.assets_dir / f"{img_id}.{image.ext}"
                    out_path.write_bytes(image.blob)
                    images.append(ImageAsset(path=out_path, page=slide_idx, id=img_id))
            if slide_text:
                text_blocks.append(TextBlock(text="\n".join(slide_text), page=slide_idx,
                                              section_path=f"Slide {slide_idx}"))

        return ParsedDoc(
            doc_id=doc_id, source_path=path, format="pptx",
            text_blocks=text_blocks, images=images, n_pages=len(prs.slides),
            meta={"parser": "python-pptx"},
        )
